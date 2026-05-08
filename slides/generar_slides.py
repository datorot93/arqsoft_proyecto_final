"""Genera dos archivos .pptx:
  - slides/experimento_asr1_asr2.pptx  (3 slides, con diagrama)
  - slides/evaluacion_asr3.pptx        (3 slides, encuesta de expertos)
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ----- Paleta de colores -----
NAVY        = RGBColor(0x1F, 0x4E, 0x78)
NAVY_DARK   = RGBColor(0x14, 0x36, 0x56)
GREEN       = RGBColor(0x2E, 0x7D, 0x32)
GREEN_LIGHT = RGBColor(0xC8, 0xE6, 0xC9)
TEAL        = RGBColor(0x00, 0x7C, 0x91)
ORANGE      = RGBColor(0xE6, 0x76, 0x22)
PURPLE      = RGBColor(0x6A, 0x1B, 0x9A)
RED         = RGBColor(0xC6, 0x28, 0x28)
LIGHT_GRAY  = RGBColor(0xEC, 0xEF, 0xF1)
LIGHT_BLUE  = RGBColor(0xE3, 0xF2, 0xFD)
LIGHT_GREEN = RGBColor(0xE8, 0xF5, 0xE9)
LIGHT_AMBER = RGBColor(0xFF, 0xF8, 0xE1)
MID_GRAY    = RGBColor(0x90, 0xA4, 0xAE)
DARK_TEXT   = RGBColor(0x21, 0x21, 0x21)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)

ROOT = Path(__file__).resolve().parent

# ====================== Helpers ======================

def new_pres():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def add_title_bar(slide, title_text, subtitle=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.85))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Inches(0.45); tf.margin_right = Inches(0.45)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = title_text
    r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = WHITE
    if subtitle:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run(); r2.text = subtitle
        r2.font.size = Pt(12); r2.font.italic = True
        r2.font.color.rgb = RGBColor(0xC8, 0xD8, 0xE8)

def add_box(slide, x, y, w, h, text, *, fill=NAVY, text_color=WHITE,
            font_size=11, bold_first=True, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
            border=None, border_pt=0.75):
    shp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if border is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = border; shp.line.width = Pt(border_pt)
    tf = shp.text_frame
    tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.03)
    tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    if isinstance(text, str): text = [text]
    for i, t in enumerate(text):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = t
        r.font.size = Pt(font_size if i == 0 else max(8, font_size - 2))
        r.font.bold = bold_first if i == 0 else False
        r.font.color.rgb = text_color
    return shp

def add_text(slide, x, y, w, h, paragraphs, *, font_size=12, color=DARK_TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, item in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if isinstance(item, str):
            text, bold, size, c = item, False, font_size, color
        else:
            text = item[0]
            bold = item[1] if len(item) > 1 else False
            size = item[2] if len(item) > 2 else font_size
            c    = item[3] if len(item) > 3 else color
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = c
    return tb

def add_arrow(slide, x1, y1, x2, y2, color=NAVY_DARK, weight=2.0):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color; conn.line.width = Pt(weight)
    ln = conn.line._get_or_add_ln()
    tail = etree.SubElement(ln, qn('a:tailEnd'))
    tail.set('type', 'triangle'); tail.set('w', 'med'); tail.set('h', 'med')
    return conn

def add_footer(slide, text):
    tb = slide.shapes.add_textbox(0, Inches(7.18), Inches(13.333), Inches(0.28))
    tf = tb.text_frame; tf.margin_left = Inches(0.45)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = text
    r.font.size = Pt(9); r.font.italic = True; r.font.color.rgb = MID_GRAY

def add_section_header(slide, x, y, w, text, color=NAVY):
    add_text(slide, x, y, w, 0.45, [(text, True, 16, color)])
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      Inches(x), Inches(y + 0.40),
                                      Inches(x + w), Inches(y + 0.40))
    line.line.color.rgb = color; line.line.width = Pt(1.5)

# ====================== EXPERIMENTO ======================

def build_experimento(prs):
    # =============== Slide 1 — Qué se busca probar ===============
    s = blank_slide(prs)
    add_title_bar(s, "Experimento de validación · ASR-1 + ASR-2",
                  "Banco Z – Línea Verde · Apertura de CDT Digital de Alto Rendimiento")

    add_section_header(s, 0.5, 1.05, 12.3, "Qué se busca probar")

    # ASR-1 columna izquierda
    add_box(s, 0.6, 1.7, 6.0, 0.55, "ASR-1  ·  Latencia",
            fill=NAVY, font_size=15)
    add_text(s, 0.75, 2.4, 5.8, 2.6, [
        ("Estímulo", True, 12),
        ("Confirmación de apertura de CDT en operación normal (800 req/h).", False, 11),
        ("", False, 5),
        ("Métrica de éxito", True, 12),
        ("P95 < 800 ms      ·      P99 < 1.500 ms (guardrail)", False, 11),
        ("", False, 5),
        ("Decisión que se valida", True, 12),
        ("Respuesta 202 Accepted tras escribir CDT + outbox.", False, 11),
        ("La reserva real en el core ocurre asíncrona — el cliente no la espera.", False, 11),
    ])

    # ASR-2 columna derecha
    add_box(s, 6.8, 1.7, 6.0, 0.55, "ASR-2  ·  Escalabilidad",
            fill=GREEN, font_size=15)
    add_text(s, 6.95, 2.4, 5.8, 2.6, [
        ("Estímulo", True, 12),
        ("Lanzamiento de tasa: 6.000 solicitudes en 20 min.", False, 11),
        ("", False, 5),
        ("Métrica de éxito", True, 12),
        ("0 pérdida  ·  P95 estable en ≥18/20 min  ·  DLQ = 0", False, 11),
        ("", False, 5),
        ("Decisión que se valida", True, 12),
        ("Kong absorbe el pico con throttling, broker desacopla del core,", False, 11),
        ("HPA escala en frío, CircuitBreaker protege contra colapso.", False, 11),
    ])

    # Caja inferior — por qué estocástico
    add_box(s, 0.6, 5.25, 12.2, 1.65, [
        "Enfoque: pruebas estocásticas, no carga determinista",
        "El ASR no se prueba con una sola medición — se prueba con la FORMA de la curva de cumplimiento bajo carga realista.",
        "Tráfico uniforme a 5 req/s entrega el volumen pero NO valida cold-start del HPA, hot-spotting por país ni resiliencia del CircuitBreaker.",
    ], fill=LIGHT_AMBER, text_color=DARK_TEXT, font_size=14, bold_first=True,
       border=ORANGE, border_pt=1.2)

    add_footer(s, "Slide 1 / 3   ·   Experimento ASR-1 / ASR-2")

    # =============== Slide 2 — Cómo se prueba ===============
    s = blank_slide(prs)
    add_title_bar(s, "Cómo se prueba",
                  "Modelo estocástico de carga + stack tecnológico + procedimiento de corrida")

    add_section_header(s, 0.5, 1.0, 6.0, "1.  Modelo estocástico de carga", color=NAVY)

    # Distribuciones
    distros = [
        ("NHPP",     "Curva del lanzamiento λ(t):\nonset abrupto → decaimiento", NAVY),
        ("MMPP-2",   "Ráfagas correlacionadas\n(viralidad, picos > 25 req/s)", TEAL),
        ("Dirichlet","Sesgo por país\nα=(3,1,1) → 60/25/15%", PURPLE),
        ("Pareto",   "Latencia core stub\n(cola pesada CICS)", ORANGE),
    ]
    for i, (name, desc, color) in enumerate(distros):
        y = 1.65 + i * 0.92
        add_box(s, 0.55, y, 1.55, 0.78, name, fill=color, font_size=14)
        add_text(s, 2.25, y + 0.05, 4.2, 0.7, [(desc, False, 10)], color=DARK_TEXT)

    add_section_header(s, 6.85, 1.0, 6.0, "2.  Stack tecnológico (paridad local ↔ OCI)", color=GREEN)
    stack_rows = [
        ("Lenguaje",     "Java 21 LTS + Spring Boot 3.3 + virtual threads"),
        ("Resiliencia",  "Resilience4j 2.2 (CircuitBreaker, Bulkhead, Retry)"),
        ("API Gateway",  "Kong Gateway 3.7 OSS  ↔  OCI API Gateway"),
        ("Mensajería",   "Redpanda 24.2 (Kafka API)  ↔  OCI Streaming"),
        ("Persistencia", "PostgreSQL 16 + CloudNativePG (1 cluster/país)"),
        ("Cluster",      "kind v0.23 + K8s 1.30  ↔  OKE 1.30"),
        ("Observabilidad","Prometheus + Grafana + Tempo + Loki + OTel"),
        ("Carga",        "k6 v0.53  ·  ramping-arrival-rate  ·  xk6-distribution"),
    ]
    for i, (cat, val) in enumerate(stack_rows):
        y = 1.65 + i * 0.45
        add_box(s, 6.95, y, 1.6, 0.36, cat, fill=GREEN, font_size=10)
        add_text(s, 8.65, y + 0.04, 4.2, 0.36, [(val, False, 10)], color=DARK_TEXT)

    # Procedimiento (bottom)
    add_section_header(s, 0.5, 5.55, 12.3, "3.  Procedimiento de corrida")
    add_box(s, 0.6, 6.05, 4.0, 0.85, [
        "Calentamiento", "5 min · 2 req/s constante · alinea JIT y pools"
    ], fill=LIGHT_GRAY, text_color=DARK_TEXT, font_size=12, border=NAVY)
    add_box(s, 4.75, 6.05, 4.0, 0.85, [
        "Línea base ASR-1", "15 min · NHPP λ ≈ 0,22 req/s · valida P95 < 800 ms"
    ], fill=LIGHT_BLUE, text_color=DARK_TEXT, font_size=12, border=NAVY)
    add_box(s, 8.9, 6.05, 4.0, 0.85, [
        "Pico ASR-2", "20 min · NHPP+MMPP-2+Dirichlet · valida 6.000 sin pérdida"
    ], fill=LIGHT_GREEN, text_color=DARK_TEXT, font_size=12, border=GREEN)

    add_footer(s, "Slide 2 / 3   ·   Repetido N=5 con seeds independientes para robustez estadística")

    # =============== Slide 3 — Diagrama del experimento ===============
    s = blank_slide(prs)
    add_title_bar(s, "Diagrama del experimento",
                  "Topología del flujo crítico de apertura de CDT  +  observabilidad transversal")

    # Layout: main flow x=2.5..10.0, sidebar x=10.5..12.9
    cx = 6.0  # column center for single-box rows
    bw_single = 3.4

    # Row 1: k6 generator (top)
    add_box(s, cx - bw_single/2, 1.05, bw_single, 0.55, [
        "k6 Generador estocástico",
    ], fill=PURPLE, font_size=12)
    add_text(s, cx - bw_single/2, 1.62, bw_single, 0.25,
             [("NHPP + MMPP-2 + Dirichlet", False, 9, MID_GRAY)],
             align=PP_ALIGN.CENTER)

    # Arrow 1
    add_arrow(s, cx, 1.92, cx, 2.20, color=NAVY_DARK)
    add_text(s, cx + 0.15, 1.95, 1.2, 0.2,
             [("HTTP / JSON", False, 8, MID_GRAY)])

    # Row 2: Kong
    add_box(s, cx - bw_single/2, 2.20, bw_single, 0.55, [
        "Kong Gateway 3.7 OSS"
    ], fill=NAVY, font_size=12)
    add_text(s, cx - bw_single/2, 2.77, bw_single, 0.25,
             [("rate-limiting · throttling · prometheus", False, 9, MID_GRAY)],
             align=PP_ALIGN.CENTER)

    # Arrow 2 (Kong -> CDT row)
    add_arrow(s, cx, 3.07, cx, 3.32, color=NAVY_DARK)

    # Row 3: CDTXPais x 3 (countries)
    cdt_w = 1.7
    cdt_xs = [cx - 2.85, cx - 0.85, cx + 1.15]  # 3 columns
    countries = ["CDTXPais · PE", "CDTXPais · MX", "CDTXPais · CO"]
    for x, name in zip(cdt_xs, countries):
        add_box(s, x, 3.32, cdt_w, 0.55, name, fill=GREEN, font_size=11)
    add_text(s, cx - 2.85, 3.89, 3*cdt_w + 1.0, 0.25,
             [("Spring Boot 3.3 · Java 21 · virtual threads · HPA min=2 max=20", False, 9, MID_GRAY)],
             align=PP_ALIGN.CENTER)

    # Arrows CDT -> DB (3 vertical)
    for x in cdt_xs:
        add_arrow(s, x + cdt_w/2, 4.18, x + cdt_w/2, 4.40, color=NAVY_DARK, weight=1.5)

    # Row 4: Postgres x 3
    pg_color = TEAL
    for x, name in zip(cdt_xs, ["Postgres PE", "Postgres MX", "Postgres CO"]):
        add_box(s, x, 4.40, cdt_w, 0.55, name, fill=pg_color, font_size=11)
    add_text(s, cx - 2.85, 4.97, 3*cdt_w + 1.0, 0.25,
             [("PostgreSQL 16 · CloudNativePG · schemas: cdt + outbox (transacción ACID conjunta)",
               False, 9, MID_GRAY)],
             align=PP_ALIGN.CENTER)

    # Merge arrows from 3 DBs to Redpanda
    rp_y = 5.42
    for x in cdt_xs:
        add_arrow(s, x + cdt_w/2, 5.18, cx, rp_y, color=NAVY_DARK, weight=1.5)

    # Row 5: Redpanda
    rp_w = 4.4
    add_box(s, cx - rp_w/2, rp_y, rp_w, 0.55,
            "Redpanda · cdt.eventos",
            fill=ORANGE, font_size=12)
    add_text(s, cx - rp_w/2, rp_y + 0.57, rp_w, 0.25,
             [("6 particiones · RF=3 · Apicurio Schema Registry", False, 9, MID_GRAY)],
             align=PP_ALIGN.CENTER)

    # Arrow 5
    add_arrow(s, cx, rp_y + 0.85, cx, rp_y + 1.05, color=NAVY_DARK)

    # Row 6: ACL
    acl_y = rp_y + 1.05
    add_box(s, cx - bw_single/2, acl_y, bw_single, 0.55,
            "ACL  ·  CircuitBreaker + Bulkhead",
            fill=RED, font_size=12)

    # ===== Sidebar: Observabilidad =====
    sb_x = 11.0; sb_w = 2.05
    add_box(s, sb_x, 1.05, sb_w, 0.5, "Observabilidad",
            fill=NAVY_DARK, font_size=12)
    obs_items = [
        ("Prometheus", "métricas + histogramas P95/P99"),
        ("Grafana",    "dashboards golden signals"),
        ("Tempo",      "trazas OTel end-to-end"),
        ("Loki",       "logs estructurados JSON"),
        ("OTel Coll.", "pipeline OTLP unificado"),
    ]
    for i, (name, sub) in enumerate(obs_items):
        y = 1.7 + i * 0.85
        add_box(s, sb_x, y, sb_w, 0.5, name,
                fill=LIGHT_BLUE, text_color=NAVY, font_size=11,
                border=NAVY, border_pt=0.5)
        add_text(s, sb_x, y + 0.52, sb_w, 0.22,
                 [(sub, False, 8, MID_GRAY)], align=PP_ALIGN.CENTER)

    # Visual divider
    div = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                 Inches(10.7), Inches(1.05),
                                 Inches(10.7), Inches(6.85))
    div.line.color.rgb = MID_GRAY; div.line.width = Pt(0.5)
    div.line.dash_style = 7  # dash

    add_footer(s, "Slide 3 / 3   ·   Línea continua = llamada síncrona   ·   Cada componente self-hosted en kind y en OKE (paridad de binario)")

# ====================== EVALUACIÓN ASR-3 ======================

def build_evaluacion(prs):
    # =============== Slide 1 — Encuesta de validación (overview) ===============
    s = blank_slide(prs)
    add_title_bar(s, "Encuesta de validación · ASR-3 Facilidad de Integración",
                  "Banco Z – Línea Verde · Cómo medimos si el diseño cumple")

    add_section_header(s, 0.5, 1.05, 12.3, "ASR-3 que se está evaluando")
    add_box(s, 0.6, 1.6, 12.2, 0.95, [
        "Línea Verde debe integrarse con sistemas externos de validación y con el core bancario",
        "para establecer la elegibilidad de Crédito Express de forma automática · meta < 10 minutos.",
    ], fill=LIGHT_BLUE, text_color=NAVY, font_size=13, bold_first=True,
       border=NAVY, border_pt=1.0)

    add_section_header(s, 0.5, 2.95, 12.3, "Cómo se mide")

    # 4 cuadros de info
    info_boxes = [
        ("10 preguntas",   "Cubren los 5 focos técnicos:\nestándares · desacople · versionado\ndocumentación · acoplamiento", NAVY),
        ("Escala 1 – 5",   "1 No cumple\n3 Cumple parcial\n5 Cumple a cabalidad", GREEN),
        ("2 evaluaciones expertas", "Aplicada a 2 expertos\nde forma independiente\npara reducir sesgo individual", TEAL),
        ("Umbral > 4.0",   "Promedio ponderado de los 2\nexpertos debe ser ESTRICTAMENTE\nmayor a 4.0 para CUMPLIDO", ORANGE),
    ]
    for i, (title, body, color) in enumerate(info_boxes):
        x = 0.6 + i * 3.07
        add_box(s, x, 3.55, 2.85, 0.55, title, fill=color, font_size=13)
        add_text(s, x + 0.05, 4.18, 2.75, 1.5, [(body, False, 11)],
                 align=PP_ALIGN.CENTER)

    # Bottom callout
    add_box(s, 0.6, 5.95, 12.2, 1.05, [
        "Lo que la encuesta NO hace",
        "No reemplaza pruebas funcionales del flujo · No mide latencia (eso lo cubre el experimento ASR-1/ASR-2) · No reabre decisiones ya cerradas con el equipo.",
    ], fill=LIGHT_AMBER, text_color=DARK_TEXT, font_size=11, bold_first=True,
       border=ORANGE, border_pt=1.0)

    add_footer(s, "Slide 1 / 3   ·   Encuesta de validación ASR-3")

    # =============== Slide 2 — Tópicos evaluados (no preguntas literales) ===============
    s = blank_slide(prs)
    add_title_bar(s, "Tópicos evaluados",
                  "10 dimensiones técnicas · ranking por peso relativo en el ASR-3")

    add_text(s, 0.5, 1.0, 12.3, 0.5,
             [("Cada tópico está anclado a componentes específicos del diagrama (ACL, MessageBroker, MotorElegibilidadXPais, etc.) — el experto debe inspeccionar la vista estructural, no opinar en abstracto.", False, 11, MID_GRAY)],
             anchor=MSO_ANCHOR.TOP)

    # 10 tópicos rankeados por peso (de más a menos importante)
    topicos = [
        ("Anti-Corruption Layer y desacople semántico", 0.15, "ACL: AdaptadorCore · TraductorDominio · CircuitBreaker", NAVY),
        ("Estándares de comunicación (OpenAPI, AsyncAPI, gRPC)", 0.12, "ApiGateway · CDTXPais · MotorElegibilidadXPais · Broker", GREEN),
        ("Asincronía y desacople temporal con el Core", 0.12, "MessageBroker · ChangeDataCapture · CoreBancoZ", TEAL),
        ("Automatización de elegibilidad bajo SLA < 10 min", 0.12, "CDT → Broker → MotorElegibilidad → CacheDistribuido", ORANGE),
        ("Versionado de eventos y schemas (Schema Registry)", 0.10, "MessageBroker · MotorElegibilidad · DetectorFraude", PURPLE),
        ("Integración con proveedores externos heterogéneos", 0.10, "ValidadorIdentidad · ConsultaProveedores · Convenios", NAVY),
        ("Versionado de APIs públicas y compatibilidad atrás", 0.08, "ApiGateway · AplicacionMovil · AplicacionWeb", GREEN),
        ("Cohesión interna y bajo acoplamiento entre subsistemas", 0.08, "LineaVerde · Integracion · Datos (bounded contexts)", TEAL),
        ("Resiliencia y aislamiento de fallos en la integración", 0.07, "CircuitBreaker · pools acotados · MessageBroker", ORANGE),
        ("Capacidad de evolución multi-país (multi-tenancy)", 0.06, "Patrón XPais · atributo pais : Pais en clases", PURPLE),
    ]
    # 2 columnas, 5 filas
    col_w = 6.05
    row_h = 0.95
    for i, (titulo, peso, comps, color) in enumerate(topicos):
        col = i % 2
        row = i // 2
        x = 0.55 + col * (col_w + 0.15)
        y = 1.65 + row * row_h

        # Barra de peso (visual)
        bar_max_w = 0.5
        bar_w = bar_max_w * (peso / 0.15)  # normalizado al máximo (15%)
        # Background bar
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(x), Inches(y),
                                Inches(0.55), Inches(0.85))
        bg.fill.solid(); bg.fill.fore_color.rgb = LIGHT_GRAY
        bg.line.fill.background()
        # Filled portion
        fill_h = 0.85 * (peso / 0.15)
        fill_y = y + (0.85 - fill_h)
        fb = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(x), Inches(fill_y),
                                Inches(0.55), Inches(fill_h))
        fb.fill.solid(); fb.fill.fore_color.rgb = color
        fb.line.fill.background()
        # Peso text
        pt = s.shapes.add_textbox(Inches(x), Inches(y + 0.85),
                                  Inches(0.55), Inches(0.2))
        tfm = pt.text_frame; tfm.margin_left = 0; tfm.margin_right = 0
        ppm = tfm.paragraphs[0]; ppm.alignment = PP_ALIGN.CENTER
        rrm = ppm.add_run(); rrm.text = f"{int(peso*100)}%"
        rrm.font.size = Pt(8); rrm.font.bold = True; rrm.font.color.rgb = color

        # Texto del tópico
        add_text(s, x + 0.65, y, col_w - 0.7, 0.85, [
            (f"{i+1:>2}.  {titulo}", True, 12, NAVY_DARK),
            (comps, False, 9, MID_GRAY),
        ], anchor=MSO_ANCHOR.TOP)

    # Leyenda inferior (5 focos cubiertos)
    add_box(s, 0.55, 6.85, 12.25, 0.45, [
        "Cobertura de focos: Estándares · Desacople de interfaces · Versionado · Documentación implícita · Acoplamiento     ·     suma de pesos = 100%",
    ], fill=LIGHT_GRAY, text_color=NAVY_DARK, font_size=10, bold_first=False,
       border=MID_GRAY, border_pt=0.5)

    add_footer(s, "Slide 2 / 3   ·   Tópicos evaluados — sin texto literal de las preguntas, solo dimensiones")

    # =============== Slide 3 — Modelo de evaluación con 2 expertos ===============
    s = blank_slide(prs)
    add_title_bar(s, "Modelo de evaluación con 2 expertos",
                  "Cómo se combinan las 2 calificaciones y cuándo el ASR se considera CUMPLIDO")

    add_section_header(s, 0.5, 1.0, 12.3, "Flujo de evaluación")

    # Diagram: 2 experts -> aggregate -> verdict
    # Box 1: Experto A
    add_box(s, 0.7, 1.7, 2.6, 0.9, [
        "Experto A", "califica las 10 preguntas\nde forma independiente"
    ], fill=NAVY, font_size=13)
    # Box 2: Experto B
    add_box(s, 0.7, 2.85, 2.6, 0.9, [
        "Experto B", "califica las 10 preguntas\nde forma independiente"
    ], fill=GREEN, font_size=13)

    # Arrow merge
    add_arrow(s, 3.4, 2.15, 4.5, 3.0, color=NAVY_DARK)
    add_arrow(s, 3.4, 3.30, 4.5, 3.0, color=NAVY_DARK)

    # Box 3: Promedio por pregunta
    add_box(s, 4.5, 2.55, 3.6, 0.95, [
        "Calificación final por pregunta",
        "sᵢ = (sᵢᴬ + sᵢᴮ) / 2"
    ], fill=TEAL, font_size=12)
    add_arrow(s, 8.1, 3.0, 8.85, 3.0, color=NAVY_DARK)

    # Box 4: Promedio ponderado
    add_box(s, 8.85, 2.55, 3.95, 0.95, [
        "Promedio ponderado global",
        "S̄ = Σ ( wᵢ · sᵢ )      con  Σ wᵢ = 1.00"
    ], fill=ORANGE, font_size=12)

    # Veredicto
    add_section_header(s, 0.5, 3.95, 12.3, "Reglas de veredicto")

    # 3 reglas
    rules = [
        ("Regla principal",
         "S̄ > 4.0  →  CUMPLIDO\nS̄ ≤ 4.0  →  NO CUMPLIDO\nUmbral estricto: 4.0 exacto NO aprueba.",
         GREEN),
        ("Regla de mínimo por dimensión",
         "Si alguna sᵢ ≤ 2,\nNO CUMPLIDO aunque S̄ > 4.0.\nUna integración con dimensión\ngravemente débil no es compensable.",
         ORANGE),
        ("Regla de discrepancia entre expertos",
         "Si |sᵢᴬ − sᵢᴮ| > 1.5\nse activa reconciliación:\nlos dos expertos se reúnen y\nentregan una calificación conjunta.",
         RED),
    ]
    for i, (title, body, color) in enumerate(rules):
        x = 0.55 + i * 4.18
        add_box(s, x, 4.55, 4.0, 0.55, title, fill=color, font_size=13)
        add_text(s, x + 0.1, 5.18, 3.85, 1.6, [(body, False, 11)],
                 align=PP_ALIGN.LEFT)

    # Salida del proceso
    add_box(s, 0.55, 6.55, 12.25, 0.55, [
        "Salida del proceso: veredicto + lista de hallazgos · cada calificación ≤ 3 → ticket de rediseño con componente específico del diagrama",
    ], fill=LIGHT_AMBER, text_color=NAVY_DARK, font_size=11, bold_first=False,
       border=ORANGE, border_pt=1.0)

    add_footer(s, "Slide 3 / 3   ·   2 expertos × 10 preguntas = 20 calificaciones agregadas en un veredicto auditable")

# ====================== Main ======================

def main():
    # Experimento
    p1 = new_pres()
    build_experimento(p1)
    out1 = ROOT / "experimento_asr1_asr2.pptx"
    p1.save(out1)
    print(f"Generado: {out1}")

    # Evaluación
    p2 = new_pres()
    build_evaluacion(p2)
    out2 = ROOT / "evaluacion_asr3.pptx"
    p2.save(out2)
    print(f"Generado: {out2}")

if __name__ == "__main__":
    main()
